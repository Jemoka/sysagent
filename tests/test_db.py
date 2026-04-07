import binascii
import csv
import json

import pytest
import yaml
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from sysagent.db import DB


def write_text_file(path, text):
    path.write_text(text, encoding="utf-8")


def write_json_file(path, _):
    path.write_text(
        json.dumps(
            {
                "title": "JSON Title",
                "tags": ["alpha", "beta"],
                "details": {"owner": "Ada"},
            }
        ),
        encoding="utf-8",
    )


def write_yaml_file(path, _):
    path.write_text(
        yaml.safe_dump(
            {
                "title": "YAML Title",
                "tags": ["gamma", "delta"],
                "details": {"owner": "Grace"},
            }
        ),
        encoding="utf-8",
    )


def write_toml_file(path, _):
    path.write_text(
        """
title = "TOML Title"
tags = ["release", "notes"]

[details]
owner = "Linus"
""".strip(),
        encoding="utf-8",
    )


def write_csv_file(path, _):
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["name", "score"])
        writer.writerow(["Ada", "10"])


def write_html_file(path, _):
    path.write_text(
        "<html><body><h1>Launch Plan</h1><p>Milestone one</p></body></html>",
        encoding="utf-8",
    )


def write_xml_file(path, _):
    path.write_text(
        "<root><title>Spec Sheet</title><item>Widget</item></root>",
        encoding="utf-8",
    )


def write_ipynb_file(path, _):
    path.write_text(
        json.dumps(
            {
                "cells": [
                    {"cell_type": "markdown", "source": ["Notebook Title\n", "Overview"]},
                    {"cell_type": "code", "source": ["print('hello world')"]},
                ],
                "metadata": {},
                "nbformat": 4,
                "nbformat_minor": 5,
            }
        ),
        encoding="utf-8",
    )


def write_docx_file(path, _):
    document = DocxDocument()
    document.add_heading("Quarterly Report")
    document.add_paragraph("Revenue increased year over year.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Sales"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "42"
    document.save(path)


def write_xlsx_file(path, _):
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Metrics"
    sheet.append(["Metric", "Value"])
    sheet.append(["Latency", 12])
    workbook.save(path)


def write_pptx_file(path, _):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Phase one complete"
    presentation.save(path)


@pytest.mark.parametrize(
    ("filename", "writer", "expected_text"),
    [
        ("notes.txt", write_text_file, "plain text document"),
        ("data.json", write_json_file, "JSON Title"),
        ("config.yaml", write_yaml_file, "YAML Title"),
        ("config.toml", write_toml_file, "TOML Title"),
        ("table.csv", write_csv_file, "Ada\t10"),
        ("page.html", write_html_file, "Launch Plan"),
        ("tree.xml", write_xml_file, "Spec Sheet"),
        ("analysis.ipynb", write_ipynb_file, "Notebook Title"),
        ("report.docx", write_docx_file, "Quarterly Report"),
        ("scores.xlsx", write_xlsx_file, "Latency\t12"),
        ("deck.pptx", write_pptx_file, "Phase one complete"),
    ],
)
def test_index_common_filetypes(tmp_path, filename, writer, expected_text):
    source = tmp_path / filename
    writer(source, expected_text)

    database = DB(tmp_path / "index")
    database.index(source, {"source": "test"})
    database.index(source, {"source": "test"})
    database.db.commit()

    file_hash = binascii.crc32(source.read_bytes()) % (1 << 32)
    cache_file = source.parent / f"{source.stem}.extracted.{file_hash}.txt"

    assert cache_file.exists()
    assert expected_text in cache_file.read_text(encoding="utf-8")
    assert database.db.get_doccount() == 1
