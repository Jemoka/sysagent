from pptx import Presentation


def _extract_pptx(path):
    presentation = Presentation(path)
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts = [f"# slide {index}"]

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append("\t".join(cells))

        slides.append("\n".join(parts))

    return "\n\n==========\n\n".join(slides)
